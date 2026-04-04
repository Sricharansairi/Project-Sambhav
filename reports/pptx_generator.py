from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def generate_pptx(data, output_path):
    """
    Professional PPTX generator for Project Sambhav.
    """
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = data.get("project_name", "Project Sambhav")
    subtitle.text = f"{data.get('subtitle', '')}\nPrediction ID: {data.get('prediction_id', 'N/A')}\n{data.get('generated_at', 'N/A')}"

    # Slide 2: Probabilistic Analysis
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Probabilistic Inference Layers"
    
    rows, cols = 4, 4
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9.0), Inches(2.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    headers = ["Layer", "Probability", "Contribution", "Status"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        
    layers = [
        ["ML Prediction Layer", data.get("ml_probability", "N/A"), "65%", "CALIBRATED"],
        ["LLM Inference Layer", data.get("llm_probability", "N/A"), "35%", "STOCHASTIC"],
        ["Reconciled Final", data.get("reconciled_probability", "N/A"), "100%", "OPTIMIZED"]
    ]
    for r, layer in enumerate(layers):
        for c, val in enumerate(layer):
            table.cell(r+1, c).text = str(val)

    # Slide 3: Outcomes
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Detailed Outcomes"
    outcomes = data.get("outcomes", [])
    if outcomes:
        rows, cols = len(outcomes) + 1, 2
        table = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.5), Inches(7.0), Inches(4.0)).table
        table.cell(0, 0).text = "Outcome Label"
        table.cell(0, 1).text = "Probability Score"
        for i, o in enumerate(outcomes):
            table.cell(i+1, 0).text = o.get("label", "")
            table.cell(i+1, 1).text = o.get("probability", "")

    # Slide 4: Input Parameters
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Input Parameters"
    params = data.get("parameters", {})
    if params:
        rows = min(len(params) + 1, 10) # Limit to 10 rows for visibility
        table = slide.shapes.add_table(rows, 2, Inches(1.0), Inches(1.5), Inches(8.0), Inches(4.5)).table
        table.cell(0, 0).text = "Parameter"
        table.cell(0, 1).text = "Value"
        for i, (k, v) in enumerate(list(params.items())[:rows-1]):
            table.cell(i+1, 0).text = k.replace("_", " ").title()
            table.cell(i+1, 1).text = str(v)

    # Slide 5: SHAP Values
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Feature Contribution (SHAP)"
    shaps = data.get("shap_values", {})
    if shaps:
        rows = min(len(shaps) + 1, 8)
        table = slide.shapes.add_table(rows, 3, Inches(0.5), Inches(1.5), Inches(9.0), Inches(4.0)).table
        table.cell(0, 0).text = "Feature"
        table.cell(0, 1).text = "Impact Score"
        table.cell(0, 2).text = "Direction"
        for i, (k, v) in enumerate(list(shaps.items())[:rows-1]):
            table.cell(i+1, 0).text = k.replace("_", " ").title()
            table.cell(i+1, 1).text = f"{v:+.4f}"
            table.cell(i+1, 2).text = "Positive (+)" if v > 0 else "Negative (-)"

    prs.save(output_path)
